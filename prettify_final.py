#!/usr/bin/env python
"""
Restyle final_presentation_slides.pptx.

Content is copied verbatim from the source deck; slide 5 (Methods), which was
an empty placeholder, is the only slide with new text. Every block is laid out
from a MEASURED wrapped height, so nothing can overlap.
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from PIL import ImageFont

SRC = "final_presentation_slides.pptx"
DST = "final_presentation_slides_pretty.pptx"
IMG = "/tmp/fps"
MEASURE_FONT = "/System/Library/Fonts/Helvetica.ttc"

INK   = RGBColor(0x0D, 0x14, 0x1C)
INK2  = RGBColor(0x3B, 0x4A, 0x59)
MUTED = RGBColor(0x6B, 0x7A, 0x87)
PAPER = RGBColor(0xF2, 0xF5, 0xF7)
CARD  = RGBColor(0xFF, 0xFF, 0xFF)
RULE  = RGBColor(0xD3, 0xDB, 0xE2)
FLAG  = RGBColor(0xB4, 0x74, 0x1F)
TEAL  = RGBColor(0x2F, 0x75, 0x65)
DISP, MONO = "Archivo", "Consolas"
LH = 1.24

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]


def lines(text, w_in, pt, bold=False):
    f = ImageFont.truetype(MEASURE_FONT, int(pt * 4))
    limit = w_in * 72 * 4
    n, cur = 0, ""
    for word in text.split():
        t = (cur + " " + word).strip()
        if f.getlength(t) * (1.06 if bold else 1.0) > limit and cur:
            n += 1; cur = word
        else:
            cur = t
    return n + (1 if cur else 0)


def height_of(text, w_in, pt, bold=False):
    return lines(text, w_in, pt, bold) * (pt / 72 * LH)


def bg(s, c=PAPER):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = c


def box(s, x, y, w, h, c, line=None):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = c
    if line: sh.line.color.rgb = line; sh.line.width = Pt(0.75)
    else: sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def txt(s, x, y, w, h, runs, *, size=14, bold=False, color=INK, font=DISP,
        align=PP_ALIGN.LEFT, caps=False, space=0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    items = [(runs, {})] if isinstance(runs, str) else runs
    p_ = tf.paragraphs[0]; p_.alignment = align; p_.line_spacing = LH
    for t, o in items:
        r = p_.add_run(); r.text = t.upper() if caps else t
        f = r.font
        f.name = o.get("font", font); f.size = Pt(o.get("size", size))
        f.bold = o.get("bold", bold); f.color.rgb = o.get("color", color)
        if o.get("space", space):
            r.font._rPr.set("spc", str(int(o.get("space", space) * 100)))
    return tb


def header(s, num, title, accent=FLAG):
    box(s, 0, 0, 0.11, 7.5, accent)
    txt(s, 0.85, 0.40, 6, 0.28, num, size=10.5, font=MONO, color=MUTED, space=1.6)
    txt(s, 0.85, 0.70, 9, 0.7, title, size=31, bold=True, color=INK)
    box(s, 0.85, 1.42, 1.6, 0.028, accent)


def bullets(s, x, y, w, items, *, l0=14.5, l1=12.5, gap0=0.13, gap1=0.055,
            budget=None):
    """Measured bullet stack. If `budget` is given, shrink to fit rather than
    overflow -- an overlap is never acceptable, a slightly smaller font is."""
    sizes = [(l0, l1)]
    if budget:
        for shrink in (1.0, 0.94, 0.88, 0.82, 0.76, 0.7):
            a, b = l0 * shrink, l1 * shrink
            tot = 0
            for lvl, t in items:
                sz = a if lvl == 0 else b
                iw = w - (0.26 if lvl == 0 else 0.60)
                tot += height_of(t, iw, sz, lvl == 0) + (gap0 if lvl == 0 else gap1)
            if tot <= budget:
                sizes = [(a, b)]; break
        else:
            sizes = [(l0 * 0.7, l1 * 0.7)]
    a, b = sizes[0]
    cy = y
    for lvl, t in items:
        sz = a if lvl == 0 else b
        iw = w - (0.26 if lvl == 0 else 0.60)
        h = height_of(t, iw, sz, lvl == 0)
        if lvl == 0:
            box(s, x, cy + sz / 72 * 0.34, 0.085, 0.085, FLAG)
            txt(s, x + 0.26, cy, iw, h, t, size=sz, bold=True, color=INK)
            cy += h + gap0
        else:
            txt(s, x + 0.32, cy, 0.2, 0.25, "–", size=sz, color=MUTED, font=MONO)
            txt(s, x + 0.60, cy, iw, h, t, size=sz, color=INK2)
            cy += h + gap1
    return cy


# ============================================================ 1  TITLE
s = prs.slides.add_slide(BLANK); bg(s)
box(s, 0, 0, 0.11, 7.5, FLAG)
txt(s, 0.9, 2.05, 11.6, 0.32, "Data Guardians",
    size=12, font=MONO, color=FLAG, space=1.8, caps=True)
txt(s, 0.9, 2.62, 11.7, 1.3,
    "Trust the data? – Agentic Quality Control for Scientific Research",
    size=38, bold=True, color=INK)
box(s, 0.9, 4.22, 2.4, 0.03, FLAG)
txt(s, 0.9, 4.48, 11.2, 0.5,
    "A Case Study of Inattentive Responding in Behavioral Science",
    size=18, color=INK2)
txt(s, 0.9, 6.42, 11.6, 0.3,
    "Shamima Hossain, Ritesh Moon, Yuheng Zhao, Alexandra Zienkiewic, Mikaela Akrenius",
    size=11.5, font=MONO, color=INK2)
txt(s, 0.9, 6.76, 11.6, 0.3, "ELLIS Summer School AI4Research",
    size=10.5, font=MONO, color=MUTED, space=1.2)

# ============================================================ 2  BACKGROUND
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "01", "Background")
bullets(s, 0.85, 1.78, 11.6, [
    (0, "In the 21st century, a plethora of behavioral sciences rely on information "
        "gathered through online experiments"),
    (0, "Benefits:"),
    (1, "fast data collection, larger samples"),
    (1, "participants from more diverse populations"),
    (0, "Main issue: data quality, leading to false conclusions"),
    (1, "Participants multitask, get distracted -> inattention to experiment/survey "
        "-> inflated correlations, false negatives and positives"),
    (1, "Participants maximize reward for completing a larger number of tasks more "
        "quickly instead of providing good quality data"),
    (1, "Participants misunderstand instructions"),
    (0, "Assessing data quality without interacting with participants is difficult"),
    (0, "Current state of the art: attention checks, task-based quality measures "
        "(e.g., accuracy), following threads on participant forums, developing "
        "statistical tools for specific experimental settings -> slow process whereas "
        "data collection needs to be fast and participant strategies change quickly"),
    (0, "Excluding participants from a study requires sturdy, pre-defined criteria"),
    (0, "We need a fast, automated, justified process for detecting participants that "
        "provided poor quality data"),
], l0=13.5, l1=12, budget=5.35)

# ============================================================ 3  TARGET PAPER
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "02", "Target Paper")
bullets(s, 0.85, 1.86, 6.5, [
    (0, "Spurious correlations between poor task performance and symptom measures"),
    (0, "Task measures: descriptive + model-based measures of performance in "
        "reinforcement learning task"),
    (0, "Symptom measures: scales assessing psychiatric symptoms (e.g., depression, "
        "anxiety, worry)"),
    (0, "Earlier work aimed to detect these based on task performance (e.g., accuracy) "
        "alone"),
    (0, "Authors introduce novel measures based on responses in survey measures "
        "(e.g., coherence)"),
], l0=13.5, budget=4.9)
box(s, 7.70, 1.86, 4.75, 2.95, CARD, line=RULE)
s.shapes.add_picture(f"{IMG}/s3_1.png", Inches(7.86), Inches(2.32), width=Inches(4.43))
txt(s, 7.70, 4.90, 4.75, 0.3, "Zorowitz et al. (2023)",
    size=9.5, font=MONO, color=MUTED, align=PP_ALIGN.CENTER, caps=True, space=1.2)

# ============================================================ 4  HYPOTHESIS
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "03", "Hypothesis and Baseline for AI")
claims = [
    "The conditions for creating spurious correlations and presence of inattentive "
    "participants are also reflected in statistical properties of the data (e.g., "
    "skewed distributions, shifted means, outliers)",
    "Could AI use baseline statistical properties of the data to trace the same "
    "participants that the researchers did using red flag heuristics?",
    "Could AI be used more generally for tracing unreliable participants and/or "
    "results based on statistical information that humans did not take a look at?",
]
y = 2.05
for i, c in enumerate(claims):
    h = height_of(c, 10.5, 15) + 0.55
    box(s, 0.85, y, 11.6, h, CARD, line=RULE)
    box(s, 0.85, y, 0.05, h, FLAG if i == 0 else TEAL)
    txt(s, 1.12, y + 0.24, 0.5, 0.3, f"0{i+1}", size=11, font=MONO,
        color=MUTED, space=1.3)
    txt(s, 1.80, y + 0.24, 10.4, h - 0.4, c, size=15, color=INK)
    y += h + 0.26

# ============================================================ 5  METHODS  (was empty)
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "04", "Methods")

txt(s, 0.85, 1.72, 5.6, 0.26, "Pipeline", size=10, font=MONO,
    color=MUTED, space=1.3, caps=True)
stages = [
    ("Prompt", "63 Spearman correlations, 7 symptom scales x 9 task measures. "
               "Two-sided, p < 0.05, deliberately uncorrected. The test prompt adds "
               "one sentence: exclude subjects that induce spurious correlations."),
    ("Agent", "Writes and executes its own Python. No fixed tool menu -- it inspects, "
              "revises and re-runs, so the path it takes is a result, not a setting."),
    ("Trace", "Nine fields per step: phase, thought, action, observation, error, "
              "revision trigger, confidence, id, timestamp."),
    ("Score", "Against a deterministic reference analysis, and against the agent's own "
              "conclusion -- did it do what it said it did?"),
]
y = 2.04
for i, (name, body_t) in enumerate(stages):
    box(s, 0.85, y, 0.30, 0.30, INK)
    txt(s, 0.85, y + 0.055, 0.30, 0.25, str(i + 1), size=12, bold=True,
        color=PAPER, font=MONO, align=PP_ALIGN.CENTER)
    txt(s, 1.30, y + 0.015, 5.0, 0.26, name, size=13.5, bold=True, color=INK)
    hb = height_of(body_t, 4.95, 10.5)
    txt(s, 1.30, y + 0.30, 4.95, hb, body_t, size=10.5, color=INK2)
    y += 0.30 + hb + 0.30
    if i < len(stages) - 1:
        txt(s, 0.97, y - 0.28, 0.3, 0.2, "↓", size=11, color=RULE, font=MONO)

txt(s, 6.95, 1.72, 5.5, 0.26, "Five conditions", size=10, font=MONO,
    color=MUTED, space=1.3, caps=True)
txt(s, 6.95, 2.02, 5.5, 0.3,
    "The prompt never changes. Only the evidence available for judging quality.",
    size=10.5, color=MUTED)
rows = [("TEST 1", "—", "no basis to exclude at all"),
        ("TEST 2", "task", "accuracy, response times"),
        ("TEST 3", "survey", "per-item responses"),
        ("TEST 4", "task + survey", "both proxies"),
        ("TEST 5", "+ metrics", "the paper's own attention checks")]
box(s, 6.95, 2.48, 5.5, 0.014, INK)
ry = 2.60
for name, ev, allows in rows:
    txt(s, 6.95, ry, 1.0, 0.26, name, size=11, bold=True, font=MONO, color=INK)
    txt(s, 8.00, ry, 1.8, 0.26, ev, size=10.5, font=MONO, color=FLAG)
    txt(s, 9.90, ry, 2.55, height_of(allows, 2.5, 10), allows, size=10, color=INK2)
    ry += 0.42
box(s, 6.95, ry - 0.05, 5.5, 0.014, RULE)
txt(s, 6.95, ry + 0.12, 5.5, 0.5,
    "TEST 1 has no basis for exclusion — declining to exclude is the correct answer.",
    size=10, color=MUTED)

fy = 6.30
box(s, 0.85, fy, 11.6, 0.86, CARD, line=RULE)
box(s, 0.85, fy, 0.05, 0.86, FLAG)
txt(s, 1.12, fy + 0.13, 11.1, 0.24, "Controls", size=10.5, bold=True, color=INK)
txt(s, 1.12, fy + 0.38, 11.1, 0.4,
    "Correlations always run on the same 386-subject table, so n is identical in every "
    "condition  ·  the paper's attention-check labels are held out and used only to "
    "score exclusions  ·  the reference analysis is deterministic (scipy), no model.",
    size=10, color=INK2)


# ============================================================ 6  AGGREGATE RESULTS
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "05", "Aggregate Results", accent=TEAL)

txt(s, 8.30, 1.72, 4.15, 0.26, "Reference — paper", size=9.5, font=MONO,
    color=MUTED, space=1.3, caps=True)
box(s, 8.30, 2.02, 4.15, 1.35, CARD, line=RULE)
s.shapes.add_picture(f"{IMG}/s6_2.png", Inches(8.38), Inches(2.09), width=Inches(3.99))

txt(s, 0.85, 1.72, 7.0, 0.26, "Ours — five data conditions", size=9.5, font=MONO,
    color=MUTED, space=1.3, caps=True)
labels = ["baseline", "test 2", "test 3", "test 4", "test 5"]
files = ["s6_3.png", "s6_4.png", "s6_5.png", "s6_6.png", "s6_7.png"]
x0, w, gap = 0.85, 2.28, 0.11
for i, (f, lab) in enumerate(zip(files, labels)):
    cx = x0 + i * (w + gap)
    box(s, cx, 3.72, w, 1.62, CARD, line=RULE)
    s.shapes.add_picture(f"{IMG}/{f}", Inches(cx + 0.05), Inches(3.77), width=Inches(w - 0.10))
    txt(s, cx, 5.42, w, 0.25, lab, size=8.5, font=MONO, color=MUTED,
        align=PP_ALIGN.CENTER, caps=True, space=1.1)

box(s, 0.85, 5.86, 11.6, 0.014, RULE)
txt(s, 0.85, 6.06, 11.6, 0.9,
    [("25 of 63 correlations are significant on all 386 subjects", {"bold": True}),
     (" — about three are expected from noise alone. Excluding the 85 participants "
      "who failed an attention check leaves ", {"color": INK2}),
     ("11", {"bold": True}),
     (". Fourteen of the original findings do not survive quality control.",
      {"color": INK2})],
    size=13, color=INK)

# ============================================================ 7  RESULTS
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "06", "Results", accent=TEAL)

txt(s, 0.85, 1.80, 5.5, 0.3, "Main results", size=16, bold=True, color=INK)
box(s, 0.85, 2.14, 1.1, 0.022, TEAL)
bullets(s, 0.85, 2.34, 5.5, [
    (1, "As expected, LLM needed more than summary information from the data to "
        "exclude participants"),
    (1, "As in paper, trial-level information from survey was more useful than "
        "trial-level information from task, and best results were achieved using both"),
    (1, "When LLM was given exclusion metrics used in the paper it made more "
        "exclusions than the paper"),
], l1=13, gap1=0.20, budget=4.4)

txt(s, 7.0, 1.80, 5.45, 0.3, "Next steps", size=16, bold=True, color=INK)
box(s, 7.0, 2.14, 1.1, 0.022, FLAG)
bullets(s, 7.0, 2.34, 5.45, [
    (1, "Ground truth – paper, AI, or something else?"),
    (1, "What kinds of processes did the LLM engage in and what information did it use"),
    (1, "Repeating process for other data reliability problems"),
    (1, "Developing tool for research and/or deriving insight for research"),
], l1=13, gap1=0.20, budget=4.4)

# ============================================================ 8  SIGNIFICANCE
s = prs.slides.add_slide(BLANK); bg(s)
header(s, "07", "Significance", accent=TEAL)
box(s, 0.85, 2.60, 11.6, 1.9, CARD, line=RULE)
box(s, 0.85, 2.60, 0.06, 1.9, TEAL)
txt(s, 1.45, 3.16, 10.6, 0.9,
    "Better participants, better science – more reliable results and applications",
    size=27, bold=True, color=INK)

prs.save(DST)
print(f"wrote {DST} — {len(prs.slides)} slides")
