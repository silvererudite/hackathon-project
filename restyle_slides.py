#!/usr/bin/env python
"""
Restyle project_slides_final_version.pptx. Content is copied verbatim -- this
changes typography, colour, and layout only.
"""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

SRC, DST = "project_slides_final_version.pptx", "project_slides_final_version_styled.pptx"
IMG = "/tmp/slideimg"

INK   = RGBColor(0x0D, 0x14, 0x1C)
INK2  = RGBColor(0x3B, 0x4A, 0x59)
MUTED = RGBColor(0x6B, 0x7A, 0x87)
PAPER = RGBColor(0xF2, 0xF5, 0xF7)
CARD  = RGBColor(0xFF, 0xFF, 0xFF)
RULE  = RGBColor(0xD3, 0xDB, 0xE2)
FLAG  = RGBColor(0xB4, 0x74, 0x1F)
TEAL  = RGBColor(0x2F, 0x75, 0x65)
DISP, MONO = "Archivo", "Consolas"

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]


def bg(s, c=PAPER):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = c


def box(s, x, y, w, h, c, line=None):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = c
    if line: sh.line.color.rgb = line; sh.line.width = Pt(0.75)
    else: sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def txt(s, x, y, w, h, runs, *, size=14, bold=False, color=INK, font=DISP,
        align=PP_ALIGN.LEFT, line=1.3, caps=False, space=0):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    items = [(runs, {})] if isinstance(runs, str) else runs
    p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = line
    for t, o in items:
        r = p.add_run(); r.text = t.upper() if caps else t
        f = r.font
        f.name = o.get("font", font); f.size = Pt(o.get("size", size))
        f.bold = o.get("bold", bold); f.color.rgb = o.get("color", color)
        if o.get("space", space):
            r.font._rPr.set("spc", str(int(o.get("space", space) * 100)))
    return tb


_FONT = "/System/Library/Fonts/Helvetica.ttc"
_LINE = 1.42          # line-height multiple, matches txt(line=1.3) rendering


def wrapped_lines(text, width_in, size_pt, bold=False):
    """How many lines this string really takes at this width and size.

    Measured with actual font metrics rather than a characters-per-line guess.
    The guess is what put the Main results column through the footer bar: it
    under-counted a 3-line bullet as 2 and the whole column drifted down.
    """
    from PIL import ImageFont
    f = ImageFont.truetype(_FONT, int(size_pt * 4))
    limit = width_in * 72 * 4
    lines, cur = 0, ""
    for word in text.split():
        trial = (cur + " " + word).strip()
        if f.getlength(trial) * (1.06 if bold else 1.0) > limit and cur:
            lines += 1
            cur = word
        else:
            cur = trial
    return lines + (1 if cur else 0)


def bullets_height(items, w, *, size=13, gap=0.055):
    """Total height these bullets will occupy -- call before laying out."""
    total = 0
    for lvl, t in items:
        sz = size if lvl == 0 else size - 1
        iw = w - (0.24 if lvl == 0 else 0.55)
        total += wrapped_lines(t, iw, sz, bold=(lvl == 0)) * (sz / 72 * _LINE) + gap
    return total


def bullets(s, x, y, w, items, *, size=13, gap=0.055):
    """items: (level, text). Level 0 gets a marker, level 1 an en dash.
    Each row is given exactly the height its text measures."""
    cy = y
    for lvl, t in items:
        sz = size if lvl == 0 else size - 1
        iw = w - (0.24 if lvl == 0 else 0.55)
        n = wrapped_lines(t, iw, sz, bold=(lvl == 0))
        h = n * (sz / 72 * _LINE)
        if lvl == 0:
            box(s, x, cy + 0.085, 0.085, 0.085, FLAG)
            txt(s, x + 0.24, cy, iw, h, t, size=sz, bold=True, color=INK)
        else:
            txt(s, x + 0.30, cy, 0.18, 0.3, "\u2013", size=sz, color=MUTED, font=MONO)
            txt(s, x + 0.55, cy, iw, h, t, size=sz, color=INK2)
        cy += h + gap
    return cy


# ===================================================== SLIDE 1
s = prs.slides.add_slide(BLANK); bg(s)
box(s, 0, 0, 0.11, 7.5, FLAG)
txt(s, 0.9, 2.28, 11.6, 0.3, "ELLIS Summer School AI4Research",
    size=11, font=MONO, color=MUTED, space=1.6, caps=True)
txt(s, 0.9, 2.82, 11.7, 1.3,
    "Trust the data? – Agentic Quality Control for Scientific Research",
    size=40, bold=True, color=INK, line=1.06)
box(s, 0.9, 4.36, 2.6, 0.03, FLAG)
txt(s, 0.9, 4.62, 11.2, 0.5,
    "A Case Study of Inattentive Responding in Behavioral Science",
    size=19, color=INK2)
txt(s, 0.9, 6.62, 11.6, 0.4,
    "Shamima Hossain, Ritesh Moon, Yuheng Zhao, Alexandra Zienkiewic, Mikaela Akrenius",
    size=11.5, font=MONO, color=MUTED)

# ===================================================== SLIDE 2
s = prs.slides.add_slide(BLANK); bg(s)
box(s, 0, 0, 0.11, 7.5, FLAG)
txt(s, 0.9, 0.62, 6, 0.3, "01", size=11, font=MONO, color=MUTED, space=1.6)
txt(s, 0.9, 0.95, 6, 0.8, "Background", size=34, bold=True, color=INK)
box(s, 0.9, 1.82, 1.8, 0.03, FLAG)

bullets(s, 0.9, 2.22, 6.6, [
    (0, "Problem: detecting spurious correlations caused by inattentive "
        "participants in online data on human subjects"),
    (0, "Target paper: Zorowitz et al. (2023)"),
    (1, "reinforcement learning task measures x psychiatric survey measures"),
    (1, "detecting and eliminating inattentive participants using task- and "
        "survey-based data quality metrics"),
    (0, "Research question: Can a LLM recognize and exclude the same participants "
        "based on statistical information alone – without information from the paper?"),
    (0, "Method: simple prompt + subsets of data"),
    (1, "measures only + trial-level task data + trial-level survey data + both + "
        "metrics from from paper"),
])

box(s, 7.85, 2.05, 4.6, 2.85, CARD, line=RULE)
s.shapes.add_picture(f"{IMG}/s2_1.png", Inches(8.0), Inches(2.5), width=Inches(4.3))
txt(s, 7.85, 4.98, 4.6, 0.3, "Zorowitz et al. (2023)",
    size=9.5, font=MONO, color=MUTED, align=PP_ALIGN.CENTER, caps=True, space=1.2)



# ===================================================== METHODS (inserted as slide 3)
def build_methods(prs, BLANK):
    s = prs.slides.add_slide(BLANK); bg(s)
    box(s, 0, 0, 0.11, 7.5, FLAG)
    txt(s, 0.9, 0.34, 6, 0.3, "02", size=11, font=MONO, color=MUTED, space=1.6)
    txt(s, 0.9, 0.64, 8, 0.7, "Methods", size=34, bold=True, color=INK)
    txt(s, 0.9, 1.30, 11.6, 0.35,
        "One prompt, one fixed correlation table, five levels of evidence. "
        "Only what the agent can learn about data quality changes.",
        size=13.5, color=INK2)

    # ---- the pipeline, left ----
    txt(s, 0.9, 1.92, 5.6, 0.25, "Pipeline", size=10, font=MONO,
        color=MUTED, space=1.3, caps=True)
    stages = [
        ("Prompt", "63 Spearman correlations, two-sided, p < 0.05, uncorrected. "
                   "Plus: exclude subjects that induce spurious correlations."),
        ("Agent", "Writes and executes its own Python. Free to inspect, revise, "
                  "and re-run; nothing about the order is fixed."),
        ("Trace", "9 fields per step: phase, thought, action, observation, error, "
                  "revision trigger, confidence, id, timestamp."),
        ("Score", "Compared against a deterministic reference analysis, and "
                  "against what the agent's own conclusion claimed."),
    ]
    y = 2.24
    for i, (name, body) in enumerate(stages):
        box(s, 0.9, y, 0.30, 0.30, INK)
        txt(s, 0.9, y + 0.055, 0.30, 0.25, str(i + 1), size=12, bold=True,
            color=PAPER, font=MONO, align=PP_ALIGN.CENTER)
        txt(s, 1.34, y + 0.01, 5.2, 0.25, name, size=13, bold=True, color=INK)
        n = wrapped_lines(body, 5.05, 10.5)
        txt(s, 1.34, y + 0.28, 5.05, n * 0.21, body, size=10.5, color=INK2, line=1.25)
        y += 0.34 + n * 0.20 + 0.20
        if i < len(stages) - 1:
            txt(s, 1.02, y - 0.30, 0.3, 0.2, "\u2193", size=11, color=RULE, font=MONO)

    # ---- conditions table, right ----
    txt(s, 7.0, 1.92, 5.45, 0.25, "Five conditions", size=10, font=MONO,
        color=MUTED, space=1.3, caps=True)
    rows = [
        ("TEST 1", "\u2014", "nothing to judge quality with"),
        ("TEST 2", "task", "accuracy, response times"),
        ("TEST 3", "survey", "per-item responses"),
        ("TEST 4", "task + survey", "both proxies"),
        ("TEST 5", "+ metrics", "the study's own attention checks"),
    ]
    hy = 2.24
    for h, x, w, al in (("", 7.0, 1.0, PP_ALIGN.LEFT),
                        ("EVIDENCE ADDED", 8.05, 1.75, PP_ALIGN.LEFT),
                        ("WHAT IT ALLOWS", 9.90, 2.55, PP_ALIGN.LEFT)):
        if h:
            txt(s, x, hy, w, 0.25, h, size=8.5, font=MONO, color=MUTED,
                space=1.1, align=al)
    box(s, 7.0, 2.50, 5.45, 0.014, INK)
    ry = 2.62
    for name, ev, allows in rows:
        txt(s, 7.0, ry, 1.0, 0.28, name, size=11, bold=True, font=MONO, color=INK)
        txt(s, 8.05, ry, 1.75, 0.28, ev, size=10.5, font=MONO, color=FLAG)
        n = wrapped_lines(allows, 2.5, 10)
        txt(s, 9.90, ry, 2.55, n * 0.2, allows, size=10, color=INK2, line=1.2)
        ry += 0.44
    box(s, 7.0, ry - 0.06, 5.45, 0.014, RULE)

    txt(s, 7.0, ry + 0.14, 5.45, 0.6,
        "Difficulty falls from TEST 1 to TEST 5. TEST 1 has no basis for exclusion "
        "at all \u2014 declining to exclude is the correct answer there.",
        size=10, color=MUTED, line=1.3)

    # ---- controls footer ----
    fy = 6.28
    box(s, 0.9, fy, 11.55, 0.9, CARD, line=RULE)
    box(s, 0.9, fy, 0.05, 0.9, FLAG)
    txt(s, 1.14, fy + 0.14, 11.1, 0.25, "Controls", size=10.5, bold=True, color=INK)
    txt(s, 1.14, fy + 0.40, 11.1, 0.4,
        "Correlations always run on the same 386-subject table, so n is identical "
        "in every condition  \u00b7  the study's own attention-check labels are held "
        "out and used only to score exclusions  \u00b7  the reference analysis is "
        "deterministic (scipy), no model involved.",
        size=10, color=INK2, line=1.3)


# ===================================================== SLIDE 3 (RESULTS)
build_methods(prs, BLANK)

s = prs.slides.add_slide(BLANK); bg(s)
box(s, 0, 0, 0.11, 7.5, TEAL)
txt(s, 0.9, 0.34, 6, 0.3, "03", size=11, font=MONO, color=MUTED, space=1.6)
txt(s, 0.9, 0.64, 6, 0.7, "Results", size=34, bold=True, color=INK)

# reference figure, right
txt(s, 8.55, 0.42, 4.2, 0.25, "Reference — paper",
    size=9, font=MONO, color=MUTED, space=1.3, caps=True)
box(s, 8.55, 0.72, 3.95, 1.28, CARD, line=RULE)
s.shapes.add_picture(f"{IMG}/s3_2.png", Inches(8.63), Inches(0.79), width=Inches(3.8))

# our five condition matrices
txt(s, 0.9, 2.16, 6, 0.25, "Ours — five data conditions",
    size=9, font=MONO, color=MUTED, space=1.3, caps=True)
labels = ["baseline", "test 2", "test 3", "test 4", "test 5"]
files = ["s3_3.png", "s3_4.png", "s3_5.png", "s3_6.png", "s3_7.png"]
x, w, gap = 0.9, 2.28, 0.11
for i, (f, lab) in enumerate(zip(files, labels)):
    cx = x + i * (w + gap)
    box(s, cx, 2.44, w, 1.56, CARD, line=RULE)
    s.shapes.add_picture(f"{IMG}/{f}", Inches(cx + 0.05), Inches(2.49), width=Inches(w - 0.10))
    txt(s, cx, 4.04, w, 0.25, lab, size=8.5, font=MONO, color=MUTED,
        align=PP_ALIGN.CENTER, caps=True, space=1.1)

box(s, 0.9, 4.38, 11.55, 0.02, RULE)

# two text columns
txt(s, 0.9, 4.58, 5.5, 0.3, "Main results", size=15, bold=True, color=INK)
bullets(s, 0.9, 4.94, 5.5, [
    (1, "As expected, LLM needed more than summary information from the data to "
        "exclude participants"),
    (1, "As in paper, trial-level information from survey was more useful than "
        "trial-level information from task, and best results were achieved using both"),
    (1, "When LLM was given exclusion metrics used in the paper it made more "
        "exclusions than the paper"),
], size=11.5, gap=0.04)

txt(s, 7.0, 4.58, 5.45, 0.3, "Next steps", size=15, bold=True, color=INK)
bullets(s, 7.0, 4.94, 5.45, [
    (1, "Ground truth – paper, AI, or something else?"),
    (1, "What kinds of processes did the LLM engage in and what information did it use"),
    (1, "Repeating process for other data reliability problems"),
    (1, "Developing tool for research and/or deriving insight for research"),
], size=11.5, gap=0.04)

box(s, 0.9, 6.92, 11.55, 0.44, CARD, line=RULE)
box(s, 0.9, 6.92, 0.05, 0.44, TEAL)
txt(s, 1.14, 7.03, 11.2, 0.3,
    [("Significance: ", {"bold": True}),
     ("increasing data quality -> more accurate results -> better real-world application", {})],
    size=12, color=INK)

prs.save(DST)
print(f"wrote {DST} — {len(prs.slides)} slides")
