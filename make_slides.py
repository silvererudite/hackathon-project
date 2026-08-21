#!/usr/bin/env python
"""Build the two-slide screening deck from the measured results."""
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

# Instrument-panel palette: cool neutrals, semantic colour carries the result.
INK      = RGBColor(0x0D, 0x14, 0x1C)
INK2     = RGBColor(0x3B, 0x4A, 0x59)
MUTED    = RGBColor(0x6B, 0x7A, 0x87)
PAPER    = RGBColor(0xF2, 0xF5, 0xF7)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
RULE     = RGBColor(0xD3, 0xDB, 0xE2)
FLAG     = RGBColor(0xB4, 0x74, 0x1F)   # flagged / caution
VERIFIED = RGBColor(0x2F, 0x75, 0x65)   # matches reference
FAIL     = RGBColor(0xB2, 0x3A, 0x42)   # silent failure
FAILBG   = RGBColor(0xF7, 0xE9, 0xEA)

DISPLAY, MONO = "Archivo", "Consolas"   # Consolas: safe cross-platform mono

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]
W = 13.333


def bg(slide, color=PAPER):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def text(slide, x, y, w, h, runs, *, size=18, bold=False, color=INK,
         font=DISPLAY, space=0, align=PP_ALIGN.LEFT, line=1.25, caps=False):
    """runs: str, or list of (text, {overrides}) for mixed formatting."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    items = [(runs, {})] if isinstance(runs, str) else runs
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line
    for t, over in items:
        r = p.add_run()
        r.text = t.upper() if caps else t
        f = r.font
        f.name = over.get("font", font)
        f.size = Pt(over.get("size", size))
        f.bold = over.get("bold", bold)
        f.color.rgb = over.get("color", color)
        if over.get("space", space):
            from pptx.oxml.ns import qn
            r.font._rPr.set("spc", str(int(over.get("space", space) * 100)))
    return box


def rule(slide, x, y, w, color=RULE, h=0.012):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    return s


def block(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background(); s.shadow.inherit = False
    return s


# =====================================================================
# SLIDE 1 — the question and the finding
# =====================================================================
s1 = prs.slides.add_slide(BLANK); bg(s1)
block(s1, 0, 0, 0.11, 7.5, FLAG)          # spine

text(s1, 0.85, 0.62, 11, 0.3,
     "ELLIS SUMMER SCHOOL 2026  ·  AI FOR RESEARCH",
     size=11, font=MONO, color=MUTED, space=1.6)

text(s1, 0.85, 1.15, 11.6, 1.5,
     "Trust the Data?", size=54, bold=True, color=INK)
text(s1, 0.85, 2.02, 11.6, 0.6,
     "Agentic Quality Control for Scientific Research",
     size=25, bold=True, color=INK2)
text(s1, 0.85, 2.58, 11.6, 0.5,
     "A Case Study of Inattentive Responding in Behavioral Science",
     size=16, color=MUTED)

rule(s1, 0.85, 3.32, 11.6)

text(s1, 0.85, 3.62, 11.6, 0.6,
     [("63 correlations across 386 participants. ", {"bold": True}),
      ("7 symptom scales against 9 task measures, "
       "p < 0.05, uncorrected — chance alone predicts about three.", {})],
     size=16, color=INK2, line=1.4)

# the 25 -> 85 -> 11 strip
stats = [(0.85, "25", "SIGNIFICANT OF 63", FLAG),
         (4.15, "85", "FAILED AN ATTENTION CHECK  (22%)", INK),
         (8.30, "11", "SURVIVE THEIR EXCLUSION", VERIFIED)]
for x, n, k, c in stats:
    text(s1, x, 4.42, 3.2, 1.0, n, size=62, bold=True, color=c, font=MONO)
    text(s1, x, 5.42, 3.3, 0.5, k, size=10.5, font=MONO, color=MUTED, space=1.2)
for x in (3.55, 7.70):
    text(s1, x, 4.62, 0.6, 0.6, "→", size=30, color=RULE, font=MONO)

block(s1, 0.85, 6.05, 11.6, 0.028, RULE)
text(s1, 0.85, 6.28, 11.6, 0.6,
     [("14 of the 25 findings were manufactured by inattentive participants.", {"bold": True}),
      ("  Which conclusions a study reaches is decided by who the analyst kept — "
       "a choice usually buried in one line of preprocessing.", {"color": INK2})],
     size=15, color=INK, line=1.4)

text(s1, 0.85, 7.02, 11.6, 0.3,
     "Shamima Hossain · Ritesh Moon · Yuheng Zhao · Alexandra Zienkiewic · Mikaela Akrenius"
     "      |      Data: Zorowitz et al., Nature Human Behaviour 2023",
     size=9.5, font=MONO, color=MUTED)

# =====================================================================
# SLIDE 2 — what the agent actually did
# =====================================================================
s2 = prs.slides.add_slide(BLANK); bg(s2)
block(s2, 0, 0, 0.11, 7.5, FAIL)

text(s2, 0.85, 0.5, 11.6, 0.3,
     "SIX CONDITIONS  ·  IDENTICAL PROMPT  ·  gpt-4.1-mini  ·  "
     "CORRELATION TABLE FIXED AT 386 SUBJECTS",
     size=10.5, font=MONO, color=MUTED, space=1.4)

text(s2, 0.85, 0.92, 11.6, 0.9,
     "Three of six conclusions described an analysis the agent never ran.",
     size=30, bold=True, color=INK)

# ---- results table ----
rows = [
    ("baseline", "—",                  "—",   "0",  "25/63", "25/63", "matches", False),
    ("test 1",   "—",                  "0",   "0",  "25/63", "25/63", "matches", False),
    ("test 2",   "task",               "171", "0",  "25/63", "13/63", "silent failure", True),
    ("test 3",   "survey",             "41",  "0",  "25/63", "23/63", "silent failure", True),
    ("test 4",   "task + survey",      "41",  "0",  "25/63", "13/63", "silent failure", True),
    ("test 5",   "+ attention checks", "85",  "85", "11/63", "11/63", "matches", False),
]
hdr = ["CONDITION", "EVIDENCE GIVEN", "FLAGGED", "ACTUALLY\nEXCLUDED",
       "AGENT\nREPORTED", "REFERENCE", ""]
colx = [0.85, 2.55, 5.05, 6.30, 7.85, 9.35, 10.85]
colw = [1.7, 2.5, 1.25, 1.55, 1.5, 1.5, 2.0]

y = 2.0
for i, (h, x, w) in enumerate(zip(hdr, colx, colw)):
    al = PP_ALIGN.LEFT if i < 2 else PP_ALIGN.RIGHT
    text(s2, x, y, w, 0.5, h, size=9, font=MONO, color=MUTED, space=1.1, align=al, line=1.15)
rule(s2, 0.85, 2.52, 11.6, INK, h=0.016)

y = 2.66
for cond, ev, fl, ex, rep, ref, tag, bad in rows:
    if bad:
        block(s2, 0.85, y - 0.045, 11.6, 0.46, FAILBG)
    vals = [(cond, INK, True), (ev, INK2, False), (fl, INK2, False),
            (ex, FAIL if bad else INK2, bad), (rep, INK, False), (ref, INK2, False)]
    for i, ((v, c, b), x, w) in enumerate(zip(vals, colx, colw)):
        al = PP_ALIGN.LEFT if i < 2 else PP_ALIGN.RIGHT
        text(s2, x, y, w, 0.35, v, size=13, font=MONO, color=c, bold=b, align=al)
    text(s2, colx[6], y + 0.015, colw[6], 0.35, tag,
         size=9.5, font=MONO, bold=True, space=1.0,
         color=FAIL if bad else VERIFIED, caps=True)
    y += 0.52
rule(s2, 0.85, y - 0.06, 11.6)

# ---- the punchline ----
py = y + 0.28
block(s2, 0.85, py, 0.045, 1.34, FAIL)
text(s2, 1.12, py + 0.02, 11.2, 0.4,
     'test 4  →  “after excluding 195 subjects…”',
     size=16, bold=True, font=MONO, color=INK)
text(s2, 1.12, py + 0.44, 11.2, 0.95,
     [("It excluded nobody.", {"bold": True}),
      ("  It flagged the subjects, wrote a conclusion describing their removal, "
       "then ran the sweep over the untouched sample. Every number it reported is "
       "internally plausible — ", {"color": INK2}),
      ("an output-only benchmark scores this as a success.", {"bold": True})],
     size=14, color=INK, line=1.35)

text(s2, 0.85, 7.02, 11.6, 0.3,
     "So we record the process, not just the answer — 9 fields per step, including what the "
     "agent did and what it believed it had done.      |      "
     "test 1 is the quiet win: asked to exclude, given nothing to exclude with, it declined.",
     size=9.5, font=MONO, color=MUTED)

prs.save("project_slides.pptx")
print("wrote project_slides.pptx —", len(prs.slides), "slides")
